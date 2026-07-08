#!/usr/bin/env python3
"""Generate Garanti Lab 1 gRPC mesh reference deck (PPTX)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("grpc-mesh-connections.pptx")

# Colors
NAVY = RGBColor(0x1A, 0x23, 0x32)
BLUE = RGBColor(0x00, 0x6E, 0xB8)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
GREEN = RGBColor(0x2B, 0x8A, 0x3E)
GRAY = RGBColor(0x5C, 0x6B, 0x7A)
DARK = RGBColor(0x2D, 0x37, 0x48)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_round_rect(slide, left, top, width, height, text, fill=LIGHT, line=BLUE, size=11, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=GRAY):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(2)
    return conn


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)
    add_textbox(
        slide, Inches(0.6), Inches(2.0), Inches(12), Inches(1),
        "gRPC in the Service Mesh", size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(0.6), Inches(2.9), Inches(12), Inches(0.6),
        "Connection pools · DestinationRule knobs · Lab 1 reference",
        size=20, color=GRAY, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(0.6), Inches(3.6), Inches(12), Inches(0.5),
        "Garanti BBVA · OSSM 3 sidecar · OpenShift",
        size=14, color=GRAY, align=PP_ALIGN.CENTER,
    )


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_textbox(
        slide, Inches(0.5), Inches(0.35), Inches(12), Inches(0.6),
        "Inside the sidecar — where TCP connections & pools are established",
        size=24, bold=True, color=BLUE,
    )

    # Fortio
    add_round_rect(slide, Inches(0.4), Inches(2.2), Inches(1.3), Inches(0.9),
                   "Fortio\n-c = parallelism", fill=RGBColor(0xE8, 0xF4, 0xFC), line=BLUE, size=10)

    # Conn A
    add_round_rect(slide, Inches(1.85), Inches(2.35), Inches(0.9), Inches(0.6),
                   "Conn A", fill=RGBColor(0xD0, 0xE8, 0xF8), line=BLUE, size=11, bold=True)
    add_arrow(slide, Inches(1.7), Inches(2.65), Inches(1.85), Inches(2.65), BLUE)

    # Client pod boundary
    pod = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.9), Inches(1.3), Inches(4.8), Inches(3.5))
    pod.fill.background()
    pod.line.color.rgb = GRAY
    pod.line.dash_style = 2
    add_textbox(slide, Inches(3.1), Inches(1.4), Inches(2), Inches(0.3), "grpc-client pod", size=10, color=GRAY)

    # Envoy box
    envoy = add_round_rect(slide, Inches(3.1), Inches(1.75), Inches(4.4), Inches(2.8),
                           "", fill=RGBColor(0xF0, 0xF4, 0xF8), line=BLUE, size=10)
    add_textbox(slide, Inches(3.2), Inches(1.85), Inches(4.2), Inches(0.35),
                "istio-proxy (Envoy)", size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    add_round_rect(slide, Inches(3.25), Inches(2.25), Inches(1.5), Inches(0.55),
                   "Inbound listener\napp → :15001", size=9)
    add_round_rect(slide, Inches(4.95), Inches(2.25), Inches(1.35), Inches(0.55),
                   "Router / LB", size=10, bold=True)

    # Outbound cluster
    cluster = add_round_rect(slide, Inches(3.15), Inches(2.95), Inches(4.3), Inches(1.45),
                             "", fill=RGBColor(0xFF, 0xF3, 0xE0), line=ORANGE, size=9)
    add_textbox(slide, Inches(3.2), Inches(3.0), Inches(4.2), Inches(0.35),
                "Outbound cluster (grpc-echo)  ◄  DestinationRule applies HERE",
                size=10, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    for i, label in enumerate(["Pool→Pod A\nTCP+HTTP/2", "Pool→Pod B\nTCP+HTTP/2", "Pool→Pod C\nTCP+HTTP/2"]):
        add_round_rect(slide, Inches(3.25 + i * 1.4), Inches(3.45), Inches(1.25), Inches(0.8),
                       label, fill=WHITE, line=ORANGE, size=8)

    add_arrow(slide, Inches(2.75), Inches(2.65), Inches(3.25), Inches(2.5), BLUE)
    add_arrow(slide, Inches(4.6), Inches(2.8), Inches(4.6), Inches(2.95), ORANGE)

    # Conn B
    add_round_rect(slide, Inches(7.85), Inches(2.35), Inches(0.9), Inches(0.6),
                   "Conn B", fill=RGBColor(0xFF, 0xE8, 0xCC), line=ORANGE, size=11, bold=True)
    add_arrow(slide, Inches(7.45), Inches(3.5), Inches(7.85), Inches(2.7), ORANGE)

    # Echo pods
    for i, y in enumerate([1.55, 2.35, 3.15]):
        add_round_rect(slide, Inches(8.95), Inches(y), Inches(1.2), Inches(0.65),
                       f"echo Pod\n{'ABC'[i]}", fill=RGBColor(0xE6, 0xF4, 0xEA), line=GREEN, size=10)
        add_arrow(slide, Inches(8.75), Inches(2.65), Inches(8.95), Inches(y + 0.32), ORANGE)

    # Legend
    leg = add_round_rect(slide, Inches(10.3), Inches(1.5), Inches(2.8), Inches(2.8),
                         "", fill=LIGHT, line=GRAY, size=9)
    lines = [
        "Established by:",
        "• Conn A — Fortio → sidecar (not DR)",
        "• Conn B — Envoy → each echo pod",
        "",
        "Pool = per backend pod IP",
        "• HTTP/2 multiplexes many RPCs",
        "• Extra -c may add pools (~6)",
        "",
        "Lab: check.sh connections",
        "(port after outbound| in logs)",
    ]
    add_textbox(slide, Inches(10.45), Inches(1.6), Inches(2.5), Inches(2.6),
                "\n".join(lines), size=9, color=DARK)

    add_textbox(
        slide, Inches(0.5), Inches(5.1), Inches(12), Inches(0.4),
        "Conn A = app→proxy  ·  Conn B + pools = trafficPolicy.connectionPool on grpc-echo host",
        size=12, bold=True, color=NAVY,
    )


def slide_table(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_textbox(
        slide, Inches(0.5), Inches(0.35), Inches(12), Inches(0.6),
        "DestinationRule knobs — layer, lab step, what to watch",
        size=24, bold=True, color=BLUE,
    )

    rows = [
        ("Setting (DR)", "Layer", "Operates on", "Lab", "Observe"),
        ("loadBalancer: ROUND_ROBIN", "HTTP/gRPC L7", "Router — which pod per RPC", "1", "watch-pods.sh · check.sh pods"),
        ("http.maxRequestsPerConnection: 1", "HTTP/2", "Recycle TCP after N RPCs", "2", "check.sh connections → churn"),
        ("tcp.maxConnections: 1", "TCP", "Cap parallel TCP per pod host", "2b", "check.sh connections (~3, -c 1)"),
        ("http.idleTimeout: 5s", "HTTP/2", "Close when no active RPCs", "3", "check.sh connections 90s"),
        ("tcp.tcpKeepalive", "TCP (kernel)", "Dead peer / NAT probes (≠ HTTP idle)", "4", "tcpdump (not port count)"),
        ("http.http2MaxRequests: 2", "HTTP/2", "Max concurrent streams per conn", "5", "check.sh errors"),
        ("Service port name grpc", "Discovery", "Protocol / cluster selection", "6", "check.sh pods"),
        ("Gateway protocol: GRPC", "Ingress", "Gateway listener (not sidecar)", "7", "test-gateway-grpc.sh"),
    ]

    tbl_shape = slide.shapes.add_table(len(rows), 5, Inches(0.4), Inches(1.1), Inches(12.5), Inches(3.8))
    table = tbl_shape.table
    col_widths = [Inches(2.8), Inches(1.1), Inches(3.2), Inches(0.55), Inches(2.4)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9 if r else 10)
                p.font.bold = r == 0
                p.font.color.rgb = WHITE if r == 0 else DARK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY

    add_textbox(
        slide, Inches(0.5), Inches(5.05), Inches(5.8), Inches(0.9),
        "Client (not DR):  Fortio -c = Conn A parallelism  ·  -qps = rate (not TCP count)",
        size=11, bold=True, color=BLUE,
    )
    add_textbox(
        slide, Inches(6.5), Inches(5.05), Inches(6), Inches(0.9),
        "Access log:  before outbound| = echo pod IP  ·  after outbound| = Conn B local port",
        size=11, bold=True, color=ORANGE,
    )


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_title(prs)
    slide_architecture(prs)
    slide_table(prs)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
